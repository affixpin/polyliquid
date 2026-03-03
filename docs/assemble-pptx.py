from pptx import Presentation
from pptx.util import Emu

prs = Presentation()
prs.slide_width = Emu(9144000)
prs.slide_height = Emu(5143500)

blank_layout = prs.slide_layouts[6]

for i in range(1, 20):
    slide = prs.slides.add_slide(blank_layout)
    img_path = f"docs/slide-images/slide-{i:02d}.png"
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)

output = "docs/polyliquid-pitch-en.pptx"
prs.save(output)
print(f"Saved PPTX with {len(prs.slides)} slides to {output}")
