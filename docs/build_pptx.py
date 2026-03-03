"""
Build Polyliquid English pitch deck matching the website's Forge theme.
Dark background (#0c0a08) + gold (#e5a823) + teal (#2dd4a8) accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Theme colors (matching website Forge theme) ──
BG          = RGBColor(0x0c, 0x0a, 0x08)  # warm black
BG_CARD     = RGBColor(0x1a, 0x17, 0x15)  # dark chocolate cards
GOLD        = RGBColor(0xe5, 0xa8, 0x23)  # brand gold
GOLD_BRIGHT = RGBColor(0xf0, 0xc0, 0x50)  # lighter gold
TEAL        = RGBColor(0x2d, 0xd4, 0xa8)  # success/accent
RED         = RGBColor(0xf0, 0x60, 0x70)  # danger
TEXT_1      = RGBColor(0xe8, 0xe6, 0xe2)  # primary text
TEXT_2      = RGBColor(0xb0, 0xad, 0xa8)  # secondary
TEXT_3      = RGBColor(0x8e, 0x8b, 0x86)  # tertiary
DIM         = RGBColor(0x6e, 0x6b, 0x66)  # dimmest
BORDER      = RGBColor(0x2c, 0x29, 0x26)  # borders
WHITE       = RGBColor(0xff, 0xff, 0xff)

# Slide dimensions (16:9)
SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)

# Margins
LEFT_M = Emu(640000)
TOP_M  = Emu(400000)
CONTENT_W = SLIDE_W - 2 * LEFT_M

FONT_HEADING = "Arial"
FONT_BODY    = "Calibri"
FONT_MONO    = "Consolas"
FONT_BRAND   = "Arial Black"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


def set_slide_bg(slide, color=BG):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=Pt(13), color=TEXT_1, bold=False, alignment=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, runs_list, alignment=PP_ALIGN.LEFT, line_spacing=None):
    """
    Add text box with multiple styled runs per paragraph.
    runs_list = [
      [ (text, font_name, font_size, color, bold), ... ],  # paragraph 1
      [ (text, font_name, font_size, color, bold), ... ],  # paragraph 2
    ]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for pi, para_runs in enumerate(runs_list):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_after = line_spacing
        for text, fname, fsize, col, bld in para_runs:
            run = p.add_run()
            run.text = text
            run.font.name = fname
            run.font.size = fsize
            run.font.color.rgb = col
            run.font.bold = bld
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=BG_CARD, border_color=BORDER):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    # Make corners more rounded
    shape.adjustments[0] = 0.05
    return shape


def add_accent_line(slide, left, top, width, color=GOLD):
    """Add a thin horizontal accent line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(27432))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, rows, cols, data, header_color=GOLD,
              col_widths=None, row_height=Emu(280000)):
    """Add a styled table."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, row_height * rows)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c] if r < len(data) and c < len(data[r]) else ""

            # Style cell
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.name = FONT_BODY

            # Cell fill
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = tcPr.makeelement(qn('a:solidFill'), {})
            if r == 0:
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '1a1715'})
                p.font.color.rgb = GOLD
                p.font.bold = True
                p.font.size = Pt(9)
            else:
                srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': '0c0a08'})
                p.font.color.rgb = TEXT_2
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

            # Borders
            for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                ln = tcPr.makeelement(qn(border_name), {'w': '6350'})
                sf = ln.makeelement(qn('a:solidFill'), {})
                sc = sf.makeelement(qn('a:srgbClr'), {'val': '2c2926'})
                sf.append(sc)
                ln.append(sf)
                tcPr.append(ln)

    return table_shape


def section_header(slide, title, subtitle=None):
    """Add standard section header with gold accent line."""
    add_accent_line(slide, LEFT_M, TOP_M, Emu(600000), GOLD)
    add_text_box(slide, LEFT_M, TOP_M + Emu(80000), CONTENT_W, Emu(500000),
                 title, FONT_BRAND, Pt(28), TEXT_1, bold=True)
    if subtitle:
        add_text_box(slide, LEFT_M, TOP_M + Emu(550000), CONTENT_W, Emu(300000),
                     subtitle, FONT_BODY, Pt(13), TEXT_3)


# ═══════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)

# Top gold line
add_accent_line(slide, Emu(0), Emu(0), SLIDE_W, GOLD)

# Brand name
add_text_box(slide, LEFT_M, Emu(1400000), CONTENT_W, Emu(700000),
             "POLYLIQUID", FONT_BRAND, Pt(56), TEXT_1, bold=True, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, LEFT_M, Emu(2100000), CONTENT_W, Emu(400000),
             "Decentralized Delegated Reputation Oracle", FONT_BODY, Pt(18),
             TEAL, alignment=PP_ALIGN.CENTER)

# Subtitle
add_text_box(slide, LEFT_M, Emu(2700000), CONTENT_W, Emu(300000),
             "Investment Presentation", FONT_BODY, Pt(14),
             TEXT_3, alignment=PP_ALIGN.CENTER)

# Footer
add_text_box(slide, LEFT_M, Emu(3200000), CONTENT_W, Emu(300000),
             "polyliquid.ai  \u2022  March 2026", FONT_BODY, Pt(12),
             TEXT_3, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "PROBLEM", "Prediction markets need oracles. Current oracles are broken.")

# Two comparison cards
card_w = Emu(3700000)
card_h = Emu(3200000)
card_y = Emu(1100000)

# UMA card (left) - red accent
left_card = add_rounded_rect(slide, LEFT_M, card_y, card_w, card_h)
add_accent_line(slide, LEFT_M + Emu(100000), card_y + Emu(80000), Emu(500000), RED)

add_text_box(slide, LEFT_M + Emu(100000), card_y + Emu(150000), card_w - Emu(200000), Emu(300000),
             "UMA (current oracle)", FONT_BODY, Pt(15), RED, bold=True)

add_rich_text(slide, LEFT_M + Emu(100000), card_y + Emu(500000), card_w - Emu(200000), Emu(2500000), [
    [("Cost of attack: ", FONT_BODY, Pt(12), TEXT_3, False), ("$750", FONT_BODY, Pt(24), RED, True)],
    [],
    [("Vote weight = token holdings", FONT_BODY, Pt(11), TEXT_3, False)],
    [],
    [("March 2025: $7M manipulation", FONT_BODY, Pt(11), TEXT_3, False)],
    [],
    [("July 2025: $160M dispute", FONT_BODY, Pt(11), TEXT_3, False)],
    [],
    [("Influence can be bought instantly", FONT_BODY, Pt(11), RED, True)],
])

# Polyliquid card (right) - teal accent
right_x = LEFT_M + card_w + Emu(200000)
right_card = add_rounded_rect(slide, right_x, card_y, card_w, card_h)
add_accent_line(slide, right_x + Emu(100000), card_y + Emu(80000), Emu(500000), TEAL)

add_text_box(slide, right_x + Emu(100000), card_y + Emu(150000), card_w - Emu(200000), Emu(300000),
             "Polyliquid", FONT_BODY, Pt(15), TEAL, bold=True)

add_rich_text(slide, right_x + Emu(100000), card_y + Emu(500000), card_w - Emu(200000), Emu(2500000), [
    [("Cost of attack: ", FONT_BODY, Pt(12), TEXT_3, False), ("$11M", FONT_BODY, Pt(24), TEAL, True)],
    [],
    [("Vote weight = stake \u00d7 reputation", FONT_BODY, Pt(11), TEXT_3, False)],
    [],
    [("Reputation is earned over months", FONT_BODY, Pt(11), TEXT_3, False)],
    [],
    [("DAO is the final arbiter", FONT_BODY, Pt(11), TEXT_3, False)],
    [],
    [("Money can be bought. Trust cannot.", FONT_BODY, Pt(11), TEAL, True)],
])

# Bottom stat
add_text_box(slide, LEFT_M, card_y + card_h + Emu(150000), CONTENT_W, Emu(300000),
             "14,700\u00d7 more expensive to attack Polyliquid", FONT_BODY, Pt(15), TEAL,
             bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 3: MARKET SIZE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "MARKET SIZE", "Prediction markets \u2014 a fast-growing segment")

# Three big metric cards
metrics = [
    ("$44B", "Prediction market\nvolume 2025"),
    ("$20B", "Addressable market\n(subjective events, conservative\n2026 forecast)"),
    ("$40M", "Annual fees\n(0.2% of volume)"),
]

card_w = Emu(2400000)
card_h = Emu(1600000)
gap = Emu(200000)
start_x = LEFT_M
y = Emu(1200000)

for i, (num, label) in enumerate(metrics):
    x = start_x + i * (card_w + gap)
    add_rounded_rect(slide, x, y, card_w, card_h)
    add_text_box(slide, x, y + Emu(250000), card_w, Emu(500000),
                 num, FONT_BRAND, Pt(36), TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Emu(100000), y + Emu(850000), card_w - Emu(200000), Emu(600000),
                 label, FONT_BODY, Pt(11), TEXT_3, alignment=PP_ALIGN.CENTER)

# Market segmentation table
table_y = Emu(3100000)
table_data = [
    ["Category", "Share", "Type", "Our market?"],
    ["Sports", "~50%", "Objective", "No"],
    ["Crypto prices", "~20%", "Chainlink", "No"],
    ["Politics, economy, culture", "~30%", "Subjective", "YES"],
]
tbl = add_table(slide, LEFT_M, table_y, CONTENT_W, 4, 4, table_data)

# Highlight "YES" cell
cell = tbl.table.cell(3, 3)
for p in cell.text_frame.paragraphs:
    p.font.color.rgb = TEAL
    p.font.bold = True


# ═══════════════════════════════════════════════════
# SLIDE 4: KEY INSIGHT
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_CARD)

add_text_box(slide, LEFT_M, Emu(600000), CONTENT_W, Emu(500000),
             "KEY INSIGHT", FONT_BRAND, Pt(28), TEXT_1, bold=True, alignment=PP_ALIGN.CENTER)

add_rich_text(slide, LEFT_M, Emu(1400000), CONTENT_W, Emu(600000), [
    [("Money can be bought instantly.", FONT_BODY, Pt(22), TEAL, False)],
    [("Trust cannot.", FONT_BODY, Pt(22), TEAL, False)],
], alignment=PP_ALIGN.CENTER)

# Formula card
card_x = LEFT_M + Emu(800000)
card_w2 = CONTENT_W - Emu(1600000)
add_rounded_rect(slide, card_x, Emu(2300000), card_w2, Emu(2200000), BG, BORDER)

add_rich_text(slide, card_x + Emu(200000), Emu(2500000), card_w2 - Emu(400000), Emu(1800000), [
    [("Influence = Stake \u00d7 Reputation", FONT_BODY, Pt(16), TEXT_1, True)],
    [],
    [("\u03ba = 1.0 \u2014 Sybil-neutral: no incentive to merge or split accounts", FONT_BODY, Pt(12), TEXT_3, False)],
    [],
    [("Reputation = volume of capital processed without errors", FONT_BODY, Pt(13), TEXT_3, False)],
    [],
    [("Security = NPV of career + LP slash + DAO", FONT_BODY, Pt(13), TEXT_3, False)],
])


# ═══════════════════════════════════════════════════
# SLIDE 5: HOW IT WORKS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "HOW IT WORKS", "Three steps: delegation \u2192 voting \u2192 dispute")

steps = [
    ("1", "LP delegates to voter", [
        "LP selects voter by track record",
        "Voter stakes LP capital on markets",
        "Free commission: 5\u201380%",
    ]),
    ("2", "Voters resolve markets", [
        "Commit-reveal (blind voting)",
        "Weight = Stake \u00d7 Rep",
        "Consensus: 50% + 1 point",
    ]),
    ("3", "Dispute \u2192 DAO", [
        "Traders dispute with 2% bond",
        "DAO votes within 48h",
        "Attacker loses everything",
    ]),
]

step_w = Emu(2400000)
step_h = Emu(2800000)
gap = Emu(200000)
y = Emu(1200000)

for i, (num, title, bullets) in enumerate(steps):
    x = LEFT_M + i * (step_w + gap)
    add_rounded_rect(slide, x, y, step_w, step_h)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Emu(100000), y + Emu(100000),
                                     Emu(400000), Emu(400000))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BG
    p.font.name = FONT_BRAND
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

    # Title
    add_text_box(slide, x + Emu(100000), y + Emu(600000), step_w - Emu(200000), Emu(300000),
                 title, FONT_BODY, Pt(13), TEXT_1, bold=True)

    # Bullets
    bullet_runs = []
    for b in bullets:
        bullet_runs.append([("\u2022  " + b, FONT_BODY, Pt(10), TEXT_3, False)])
        bullet_runs.append([])
    add_rich_text(slide, x + Emu(100000), y + Emu(1000000), step_w - Emu(200000), Emu(1500000), bullet_runs)


# ═══════════════════════════════════════════════════
# SLIDE 6: REPUTATION FORMULA
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "REPUTATION FORMULA", "Each correct vote adds stake. Old votes decay.")

# Formula box
add_rounded_rect(slide, LEFT_M, Emu(1100000), CONTENT_W, Emu(500000), BG_CARD, TEAL)
add_text_box(slide, LEFT_M + Emu(150000), Emu(1200000), CONTENT_W - Emu(300000), Emu(300000),
             "R(t) = \u03a3 LP_stake \u00d7 0.995 ^ (days since vote)          Ceiling = stake_day / 0.005",
             FONT_MONO, Pt(13), TEAL)

# Explanation bullets
explanations = [
    ("Each correct vote adds the LP stake amount to reputation.", TEXT_1, False),
    ("Old votes decay: \u22120.5% per day. Half-life: 138 days.", TEXT_1, False),
    ("Ceiling at $1M/day stake: $1M / 0.005 = $200M.", TEAL, True),
    ("If a voter stops voting \u2014 reputation halves in 138 days.", TEXT_3, False),
    ("Protection against abandoned accounts and dormant voters.", TEXT_3, False),
]

y = Emu(1900000)
for text, color, bold in explanations:
    add_text_box(slide, LEFT_M + Emu(100000), y, CONTENT_W - Emu(200000), Emu(250000),
                 text, FONT_BODY, Pt(13), color, bold)
    y += Emu(320000)

# Bottom callout
add_rounded_rect(slide, LEFT_M, y + Emu(200000), CONTENT_W, Emu(400000), BG_CARD, TEAL)
add_text_box(slide, LEFT_M + Emu(150000), y + Emu(270000), CONTENT_W - Emu(300000), Emu(300000),
             "Reputation is the volume of capital processed without errors. It cannot be bought \u2014 only earned.",
             FONT_BODY, Pt(12), TEAL, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 7: VOTE WEIGHT & SYBIL PROTECTION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "VOTE WEIGHT & SYBIL PROTECTION",
               "\u03ba = 1.0: Sybil-neutral. No single group dominates.")

# Formula box
add_rounded_rect(slide, LEFT_M, Emu(1100000), CONTENT_W, Emu(450000), BG_CARD, TEAL)
add_text_box(slide, LEFT_M + Emu(150000), Emu(1200000), CONTENT_W - Emu(300000), Emu(300000),
             "W = LP_stake \u00d7 Rep          \u03ba = 1.0 \u2014 linear, Sybil-neutral",
             FONT_MONO, Pt(13), TEAL)

# Power distribution table
table_data = [
    ["Voter", "Count", "LP Stake", "Reputation", "Weight Share"],
    ["Elite", "10", "$1M", "$200M", "30%"],
    ["Top", "40", "$500K", "$100M", "30%"],
    ["Medium", "150", "$500K", "$30M", "34%"],
    ["Newcomer", "800", "$500K", "$1M", "6%"],
    ["Whale (no rep)", "\u2014", "$10M", "$0", "0%"],
]
add_table(slide, LEFT_M, Emu(1800000), CONTENT_W, 6, 5, table_data)

# Bottom callout
add_rounded_rect(slide, LEFT_M, Emu(3800000), CONTENT_W, Emu(400000), BG_CARD, TEAL)
add_text_box(slide, LEFT_M + Emu(150000), Emu(3870000), CONTENT_W - Emu(300000), Emu(300000),
             "30/30/34 \u2014 no single group controls the system. A whale with $10M but no rep = 0% influence.",
             FONT_BODY, Pt(12), TEAL, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 8: REPUTATION SLASHING
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "REPUTATION SLASHING",
               "Wrong vote \u2014 lose reputation. Lost dispute \u2014 lose everything.")

# Formula
add_rounded_rect(slide, LEFT_M, Emu(1100000), CONTENT_W, Emu(450000), BG_CARD, TEAL)
add_text_box(slide, LEFT_M + Emu(150000), Emu(1200000), CONTENT_W - Emu(300000), Emu(300000),
             "S_rep = 50% \u00d7 (W_majority \u2212 50%) / 50%          R_new = R \u00d7 (1 \u2212 S_rep)",
             FONT_MONO, Pt(12), TEAL)

# Example table
table_data = [
    ["Voting", "Slash", "Rep After", "Lost", "Recovery"],
    ["51/49", "1%", "$198M", "$2M", "~2 days"],
    ["60/40", "10%", "$180M", "$20M", "~20 days"],
    ["80/20", "30%", "$140M", "$60M", "~60 days"],
    ["Dispute", "100%", "$0", "$200M", "~12 months"],
]
add_table(slide, LEFT_M, Emu(1800000), CONTENT_W, 5, 5, table_data)

# Cascade warning
add_text_box(slide, LEFT_M + Emu(100000), Emu(3600000), CONTENT_W - Emu(200000), Emu(300000),
             "Cascade: rep dropped \u2192 LP slash rises \u2192 LP withdraws \u2192 stake drops \u2192 income drops",
             FONT_BODY, Pt(11), RED)

add_text_box(slide, LEFT_M + Emu(100000), Emu(3900000), CONTENT_W - Emu(200000), Emu(300000),
             "Example: elite voter, rep $200M", FONT_BODY, Pt(12), TEXT_1, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 9: LP SLASHING
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "LP SLASHING", "Two modes: regular error and dispute")

# Formula cards side by side
f_w = Emu(3700000)
f_h = Emu(700000)
f_y = Emu(1200000)

# Regular error
add_rounded_rect(slide, LEFT_M, f_y, f_w, f_h, BG_CARD, TEAL)
add_rich_text(slide, LEFT_M + Emu(150000), f_y + Emu(100000), f_w - Emu(300000), Emu(500000), [
    [("Regular error:", FONT_MONO, Pt(11), TEAL, False)],
    [("LP_slash = 20% / (1 + R / $33.3M)\u00b2", FONT_MONO, Pt(11), TEAL, False)],
])

# Dispute
add_rounded_rect(slide, LEFT_M + f_w + Emu(200000), f_y, f_w, f_h, BG_CARD, RED)
add_rich_text(slide, LEFT_M + f_w + Emu(350000), f_y + Emu(100000), f_w - Emu(300000), Emu(500000), [
    [("On dispute:", FONT_MONO, Pt(11), RED, False)],
    [("LP_slash = 20% (fixed)", FONT_MONO, Pt(11), RED, False)],
])

# LP slash table
table_data = [
    ["Voter Rep", "Regular Error", "On Dispute"],
    ["$0 (newcomer)", "20%", "20%"],
    ["$30M (medium)", "5.2%", "20%"],
    ["$100M (top)", "1.25%", "20%"],
    ["$200M (elite)", "0.5%", "20%"],
]
add_table(slide, LEFT_M, Emu(2200000), CONTENT_W, 5, 3, table_data)


# ═══════════════════════════════════════════════════
# SLIDE 10: NPV HOW WE CALCULATE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "NPV: HOW WE CALCULATE",
               "What happens to an elite voter after an attack. Step-by-step breakdown.")

# Initial data card
add_rounded_rect(slide, LEFT_M, Emu(1100000), CONTENT_W, Emu(1700000))
add_text_box(slide, LEFT_M + Emu(150000), Emu(1150000), CONTENT_W - Emu(300000), Emu(250000),
             "Elite voter initial data", FONT_BODY, Pt(14), TEXT_1, bold=True)

add_rich_text(slide, LEFT_M + Emu(150000), Emu(1450000), CONTENT_W - Emu(300000), Emu(1200000), [
    [("LP stake: $1M  \u2022  Reputation: $200M  \u2022  \u03ba = 1.0", FONT_BODY, Pt(12), TEXT_1, True)],
    [("Elite share of pool: 30% (10 elite \u00d7 $1M \u00d7 $200M = 30% of total weight)", FONT_BODY, Pt(11), TEXT_3, False)],
    [("Pool rewards: $32M \u00d7 30% / 10 elite = $960K/year", FONT_BODY, Pt(11), TEXT_3, False)],
    [("Voter commission: 80% \u2192 income = $768K/year", FONT_BODY, Pt(12), TEAL, True)],
    [("LP receives 20% = $192K \u2192 APY = 19.2% on $1M stake", FONT_BODY, Pt(11), TEXT_3, False)],
])

# After attack card
add_rounded_rect(slide, LEFT_M, Emu(3000000), CONTENT_W, Emu(1800000), BG_CARD, RED)
add_text_box(slide, LEFT_M + Emu(150000), Emu(3050000), CONTENT_W - Emu(300000), Emu(250000),
             "After attack: cascade of destruction", FONT_BODY, Pt(14), RED, bold=True)

cascade_steps = [
    ("1. DAO overturns R1 result \u2192 attackers' rep zeroed", TEXT_1, False),
    ("2. Rep = 0 \u2192 LP slash rises from 0.5% to 20%", TEXT_1, False),
    ("3. LP sees 20% slash risk \u2192 immediately withdraws capital", TEXT_1, False),
    ("4. No LP = no stake = no weight = no rewards = no income", RED, True),
    ("5. Voter starts from scratch as newcomer. Recovery: ~12 months.", TEXT_1, False),
]

y = Emu(3350000)
for text, color, bold in cascade_steps:
    add_text_box(slide, LEFT_M + Emu(150000), y, CONTENT_W - Emu(300000), Emu(220000),
                 text, FONT_BODY, Pt(11), color, bold)
    y += Emu(260000)


# ═══════════════════════════════════════════════════
# SLIDE 11: NPV MONTHLY BREAKDOWN
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "NPV: MONTHLY BREAKDOWN",
               "Elite. Income $768K/year. Rep \u2192 0 after attack.")

table_data = [
    ["Period", "Normal", "Actual", "Loss", "What happens"],
    ["Month 1\u20133", "$192K", "$600", "$191K", "Rep = 0. LP left. Minimal stake. Zero weight."],
    ["Month 3\u20136", "$192K", "$5K", "$187K", "Small LP risked $50K. Growing tier. Slash still 12%."],
    ["Month 6\u20139", "$192K", "$18K", "$174K", "Medium tier. $200K stake. Rep ~$30M. Slash 5%."],
    ["Month 9\u201312", "$192K", "$96K", "$96K", "Near top tier. $500K stake. Rep ~$80M."],
    ["TOTAL", "$768K", "$120K", "$653K", "NPV of elite voter = $653K"],
]
tbl = add_table(slide, LEFT_M, Emu(1100000), CONTENT_W, 6, 5, table_data,
                col_widths=[Emu(1200000), Emu(1100000), Emu(1100000), Emu(1100000), Emu(3400000)])

# Highlight total row
for c in range(5):
    cell = tbl.table.cell(5, c)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = GOLD
        p.font.bold = True

# Note
add_rounded_rect(slide, LEFT_M, Emu(3600000), CONTENT_W, Emu(800000))
add_text_box(slide, LEFT_M + Emu(150000), Emu(3650000), CONTENT_W - Emu(300000), Emu(250000),
             "Why NPV is a conservative estimate", FONT_BODY, Pt(13), TEXT_1, bold=True)
add_rich_text(slide, LEFT_M + Emu(150000), Emu(3950000), CONTENT_W - Emu(300000), Emu(400000), [
    [("We assume LP returns after 12 months. In reality:", FONT_BODY, Pt(11), TEXT_3, False)],
    [("\u2022 Other voters will take their share of the pool.", FONT_BODY, Pt(11), TEXT_1, False)],
    [("\u2022 Real loss \u2248 $768K+ (entire career), not $653K.", FONT_BODY, Pt(11), RED, True)],
])


# ═══════════════════════════════════════════════════
# SLIDE 12: COST OF ATTACK
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "COST OF ATTACK",
               "Worst case: 5 elite + 20 top accounts")

# NPV by tier table
add_text_box(slide, LEFT_M, Emu(1100000), Emu(3500000), Emu(250000),
             "NPV by tier", FONT_BODY, Pt(14), TEXT_1, bold=True)

table_data_1 = [
    ["Tier", "Income/year", "NPV"],
    ["Elite (80%)", "$768K", "$653K"],
    ["Top (60%)", "$144K", "$122K"],
    ["Medium (40%)", "$40K", "$18K"],
]
add_table(slide, LEFT_M, Emu(1400000), Emu(3500000), 4, 3, table_data_1)

# Attack composition table
add_text_box(slide, LEFT_M + Emu(3800000), Emu(1100000), Emu(3500000), Emu(250000),
             "Attack: 5 elite + 20 top", FONT_BODY, Pt(14), TEXT_1, bold=True)

table_data_2 = [
    ["Component", "Amount"],
    ["NPV: 5\u00d7$653K + 20\u00d7$122K", "$5.71M"],
    ["LP slash: 5\u00d7$1M + 20\u00d7$500K \u00d720%", "$3M"],
    ["DAO bond 2%", "$2M"],
    ["TOTAL", "$10.71M"],
]
tbl2 = add_table(slide, LEFT_M + Emu(3800000), Emu(1400000), Emu(3700000), 5, 2, table_data_2)

# Highlight total
cell = tbl2.table.cell(4, 0)
for p in cell.text_frame.paragraphs:
    p.font.color.rgb = GOLD
    p.font.bold = True
cell = tbl2.table.cell(4, 1)
for p in cell.text_frame.paragraphs:
    p.font.color.rgb = GOLD
    p.font.bold = True

# Big metrics
y_metrics = Emu(3100000)
add_rounded_rect(slide, LEFT_M, y_metrics, Emu(3500000), Emu(1000000))
add_text_box(slide, LEFT_M + Emu(100000), y_metrics + Emu(100000), Emu(3300000), Emu(200000),
             "R1 protects up to", FONT_BODY, Pt(11), TEXT_3)
add_text_box(slide, LEFT_M + Emu(100000), y_metrics + Emu(350000), Emu(3300000), Emu(400000),
             "$11M (99.5% of markets)", FONT_BODY, Pt(18), TEAL, bold=True)

add_rounded_rect(slide, LEFT_M + Emu(3800000), y_metrics, Emu(3700000), Emu(1000000))
add_text_box(slide, LEFT_M + Emu(3900000), y_metrics + Emu(100000), Emu(3500000), Emu(200000),
             "With DAO", FONT_BODY, Pt(11), TEXT_3)
add_text_box(slide, LEFT_M + Emu(3900000), y_metrics + Emu(350000), Emu(3500000), Emu(400000),
             "\u221e (any size)", FONT_BODY, Pt(18), TEAL, bold=True)

add_text_box(slide, LEFT_M, Emu(4300000), CONTENT_W, Emu(300000),
             "Real attack cost is higher: coordinating 25 people, one refusal = exposure for all",
             FONT_BODY, Pt(11), TEXT_3)


# ═══════════════════════════════════════════════════
# SLIDE 13: DISPUTE TIMELINE
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "DISPUTE: TIMELINE", "R1 \u2192 Dispute window \u2192 DAO")

table_data = [
    ["Phase", "< $10M", "> $10M"],
    ["Commit", "2 hours", "2 hours"],
    ["Reveal", "2 hours", "2 hours"],
    ["Dispute window", "4 hours", "20 hours"],
    ["Total R1", "8 hours", "24 hours"],
    ["+ DAO", "+ 48h = 56h", "+ 48h = 72h"],
]
add_table(slide, LEFT_M, Emu(1200000), CONTENT_W, 6, 3, table_data)

# Big stat
add_rounded_rect(slide, LEFT_M, Emu(3400000), CONTENT_W, Emu(600000), BG_CARD, TEAL)
add_text_box(slide, LEFT_M, Emu(3500000), CONTENT_W, Emu(400000),
             "99.5% of markets are resolved within 8 hours", FONT_BODY, Pt(18), TEAL,
             bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 14: SLASH DISTRIBUTION
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "SLASH DISTRIBUTION",
               "Correct side gets rewarded. Polyliquid DAO earns from security.")

table_data = [
    ["Scenario", "R1 Slashes", "Bond", "Client\n(Polymarket)", "DAO\nPolyliquid"],
    ["No dispute", "\u2192 correct voters", "\u2014", "$0", "$0"],
    ["DAO = R1\n(spam dispute)", "75% \u2192 correct R1", "100% \u2192 client", "100% of bond", "25% of slashes"],
    ["DAO \u2260 R1\n(attack)", "50% \u2192 correct R1\n25% \u2192 client\n25% \u2192 DAO", "Returned\nto disputers", "25% of slashes", "25% of slashes"],
]
add_table(slide, LEFT_M, Emu(1200000), CONTENT_W, 4, 5, table_data,
          col_widths=[Emu(1300000), Emu(2000000), Emu(1500000), Emu(1500000), Emu(1600000)],
          row_height=Emu(450000))

# Revenue examples
y = Emu(3400000)
card_w = Emu(3700000)
add_rounded_rect(slide, LEFT_M, y, card_w, Emu(700000))
add_text_box(slide, LEFT_M + Emu(150000), y + Emu(50000), card_w - Emu(300000), Emu(200000),
             "Client revenue (attack on $100M market)", FONT_BODY, Pt(12), TEXT_1, bold=True)
add_rich_text(slide, LEFT_M + Emu(150000), y + Emu(300000), card_w - Emu(300000), Emu(350000), [
    [("Spam: 75% of bond ($1.5M)", FONT_BODY, Pt(11), TEXT_3, False)],
    [("Attack: 25% of slashes ($500K)", FONT_BODY, Pt(11), TEXT_3, False)],
])

add_rounded_rect(slide, LEFT_M + card_w + Emu(200000), y, card_w, Emu(700000))
add_text_box(slide, LEFT_M + card_w + Emu(350000), y + Emu(50000), card_w - Emu(300000), Emu(200000),
             "DAO Polyliquid revenue (attack on $100M)", FONT_BODY, Pt(12), TEXT_1, bold=True)
add_rich_text(slide, LEFT_M + card_w + Emu(350000), y + Emu(300000), card_w - Emu(300000), Emu(350000), [
    [("Spam: 25% of slashes ($500K)", FONT_BODY, Pt(11), TEXT_3, False)],
    [("Attack: 25% of slashes ($500K)", FONT_BODY, Pt(11), TEXT_3, False)],
])


# ═══════════════════════════════════════════════════
# SLIDE 15: ECONOMICS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "ECONOMICS", "Voter commissions and LP returns")

table_data = [
    ["Tier", "LP Stake", "LP Slash", "Commission", "APY LP", "Voter Income"],
    ["Elite", "$1M", "0.5%", "80%", "15\u201320%", "$768K/year"],
    ["Top", "$500K", "1.25%", "60%", "25\u201335%", "$144K/year"],
    ["Medium", "$500K", "5.2%", "40%", "50\u201370%", "$40K/year"],
    ["Growing", "$500K", "12.4%", "20%", "70\u2013100%", "$10K/year"],
    ["Newcomer", "$500K", "20%", "5%", "100\u2013200%", "$1.6K/year"],
]
add_table(slide, LEFT_M, Emu(1200000), CONTENT_W, 6, 6, table_data)


# ═══════════════════════════════════════════════════
# SLIDE 16: COMPETITIVE COMPARISON
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "COMPETITIVE COMPARISON")

table_data = [
    ["", "UMA", "Polyliquid"],
    ["Type", "Decentralized", "Decentralized"],
    ["Cost of attack", "$750", "$11M"],
    ["Vote weight", "Tokens ($)", "Stake \u00d7 Rep (\u03ba=1)"],
    ["Decentralization", "Whale dominates", "30/30/34 balance"],
    ["Cycle time", "~48h", "8\u201324h (99%)"],
    ["2025 exploits", "$7M + $160M", "\u221e with DAO"],
    ["Client revenue", "$0", "25% slashes + 75% bonds"],
]
tbl = add_table(slide, LEFT_M, Emu(1100000), CONTENT_W, 8, 3, table_data,
                col_widths=[Emu(2200000), Emu(2600000), Emu(3100000)])

# Highlight Polyliquid column values
for r in range(2, 8):
    cell = tbl.table.cell(r, 2)
    for p in cell.text_frame.paragraphs:
        p.font.color.rgb = TEAL

# Highlight UMA attack cost in red
cell = tbl.table.cell(2, 1)
for p in cell.text_frame.paragraphs:
    p.font.color.rgb = RED


# ═══════════════════════════════════════════════════
# SLIDE 17: ALL PARAMETERS
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "ALL PARAMETERS", "11 parameters. Simple and clear.")

table_data = [
    ["Parameter", "Value", "Function"],
    ["\u03ba", "1.0", "Linear weight, Sybil-neutral"],
    ["\u03b4", "0.995/day", "Rep decay (half-life 138 days)"],
    ["S_base", "50%", "Max reputation slash"],
    ["LP_slash_max", "20%", "Max LP slash"],
    ["LP_slash_dispute", "20%", "Fixed slash on dispute"],
    ["\u03c0", "20%", "Protocol share"],
    ["Fee", "0.2%", "Fixed on volume"],
    ["Bond", "2% of market", "Threshold for DAO"],
    ["MaxPayout", "99.9%", "For security calculations"],
    ["Dispute window", "4h / 20h", "< $10M / > $10M"],
]
add_table(slide, LEFT_M, Emu(1100000), CONTENT_W, 11, 3, table_data,
          col_widths=[Emu(2200000), Emu(2000000), Emu(3700000)])


# ═══════════════════════════════════════════════════
# SLIDE 18: ROADMAP
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide)
section_header(slide, "ROADMAP", "From UMA aggregator to Polymarket oracle")

phases = [
    ("PHASE 1", "UMA Voter Hub", "Q2 2026", [
        "Voter site for UMA.",
        "Voters route votes through us.",
        "Pay 2x market rewards.",
        "Build community & track record.",
    ]),
    ("PHASE 2", "Limitless (Base)", "Q2 2026", [
        "Connect Limitless on Base.",
        "Subjective events.",
        "First live client.",
        "Agreement in place.",
    ]),
    ("PHASE 3", "Prediction Markets", "Q3 2026", [
        "Scale: Probabal, Opinion",
        "and other prediction markets.",
        "Prove reliability on",
        "live markets.",
    ]),
    ("PHASE 4", "Polymarket", "Q4 2026", [
        "Propose to replace UMA.",
        "Motivation: security.",
        "$750 \u2192 $11M attack cost.",
        "50% slashes \u2192 Polymarket.",
    ]),
]

phase_w = Emu(1800000)
phase_h = Emu(2800000)
gap = Emu(120000)
y = Emu(1200000)

for i, (phase_label, title, quarter, bullets) in enumerate(phases):
    x = LEFT_M + i * (phase_w + gap)
    add_rounded_rect(slide, x, y, phase_w, phase_h)

    # Phase label
    add_text_box(slide, x + Emu(100000), y + Emu(80000), phase_w - Emu(200000), Emu(200000),
                 phase_label, FONT_BODY, Pt(10), GOLD, bold=True)

    # Title
    add_text_box(slide, x + Emu(100000), y + Emu(350000), phase_w - Emu(200000), Emu(300000),
                 title, FONT_BODY, Pt(13), TEXT_1, bold=True)

    # Quarter
    add_text_box(slide, x + Emu(100000), y + Emu(650000), phase_w - Emu(200000), Emu(200000),
                 quarter, FONT_BODY, Pt(10), TEXT_3)

    # Gold separator
    add_accent_line(slide, x + Emu(100000), y + Emu(900000), Emu(400000), GOLD)

    # Bullets
    bullet_text = "\n".join(bullets)
    add_text_box(slide, x + Emu(100000), y + Emu(1050000), phase_w - Emu(200000), Emu(1600000),
                 bullet_text, FONT_BODY, Pt(9), TEXT_3)

# Timeline line connecting phases
line_y = y + Emu(200000)
add_accent_line(slide, LEFT_M, y - Emu(40000), CONTENT_W, GOLD)


# ═══════════════════════════════════════════════════
# SLIDE 19: CLOSING
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_CARD)

# Top gold line
add_accent_line(slide, Emu(0), Emu(0), SLIDE_W, GOLD)

# Brand
add_text_box(slide, LEFT_M, Emu(1100000), CONTENT_W, Emu(700000),
             "POLYLIQUID", FONT_BRAND, Pt(52), TEXT_1, bold=True, alignment=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide, LEFT_M, Emu(1800000), CONTENT_W, Emu(400000),
             "Money can be bought. Trust cannot.", FONT_BODY, Pt(20), TEAL,
             alignment=PP_ALIGN.CENTER)

# Stats
add_text_box(slide, LEFT_M, Emu(2600000), CONTENT_W, Emu(300000),
             "polyliquid.ai", FONT_BODY, Pt(16), TEXT_3, alignment=PP_ALIGN.CENTER)

add_text_box(slide, LEFT_M, Emu(3000000), CONTENT_W, Emu(300000),
             "7 formulas  \u2022  11 parameters  \u2022  R1 + DAO", FONT_BODY, Pt(12), TEXT_3,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, LEFT_M, Emu(3400000), CONTENT_W, Emu(300000),
             "$11M attack cost  \u2022  $750 for UMA  \u2022  14,700\u00d7 difference", FONT_BODY, Pt(12), TEAL,
             alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = "docs/polyliquid-pitch-en.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
